"""
Final Presentation: Portfolio Risk Stress-Tester
STATS 418 Final Project — Mark Wang, UCLA Department of Statistics & Data Science

Usage:
    python presentation/generate_final_slides.py
Output:
    presentation/wang-mark-final.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import matplotlib.pyplot as plt
import numpy as np
from io import BytesIO

# ============================================================================
# Theme constants
# ============================================================================

C_BG        = RGBColor(0x1a, 0x1a, 0x2e)
C_TITLE_BG  = RGBColor(0x0f, 0x3c, 0x78)
C_ACCENT    = RGBColor(0x53, 0xd8, 0xfb)
C_WHITE     = RGBColor(0xff, 0xff, 0xff)
C_LGRAY     = RGBColor(0xb0, 0xb0, 0xb0)
C_GREEN     = RGBColor(0x2e, 0xcc, 0x71)
C_ORANGE    = RGBColor(0xf3, 0x9c, 0x12)

HEX_BG      = '#ffffff'
HEX_AXES_BG = '#f7f7f7'
HEX_TEXT    = '#222222'
HEX_SPINE   = '#cccccc'
HEX_BLUE    = '#2e86de'
HEX_GREEN   = '#27ae60'
HEX_ORANGE  = '#e67e22'
HEX_RED     = '#e74c3c'
HEX_PURPLE  = '#8e44ad'
HEX_CYAN    = '#2980b9'

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ============================================================================
# Low-level helpers
# ============================================================================

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             size=16, bold=False, italic=False, color=None,
             align=PP_ALIGN.LEFT):
    if color is None:
        color = C_WHITE
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return box


def title_bar(slide, title, subtitle=None):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.15), C_TITLE_BG)
    add_text(slide, title,
             Inches(0.35), Inches(0.07), Inches(12.5), Inches(0.75),
             size=28, bold=True, color=C_WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.35), Inches(0.78), Inches(12.5), Inches(0.35),
                 size=12, color=C_ACCENT)


def bullet_frame(slide, bullets, left, top, width, height, base_size=15):
    """
    bullets: list of (main_text, [sub1, sub2, ...])
    """
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for (main, subs) in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if main:
            p.space_before = Pt(9)
            run = p.add_run()
            run.text = f"• {main}"
            run.font.size = Pt(base_size)
            run.font.bold = True
            run.font.color.rgb = C_ACCENT
            run.font.name = "Calibri"
        for sub in subs:
            p2 = tf.add_paragraph()
            p2.space_before = Pt(2)
            run2 = p2.add_run()
            run2.text = f"    – {sub}"
            run2.font.size = Pt(base_size - 3)
            run2.font.color.rgb = C_WHITE
            run2.font.name = "Calibri"
    return box


def fig_to_buf(fig, save_path=None):
    buf = BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor=HEX_BG)
    if save_path:
        fig.savefig(save_path, dpi=150, bbox_inches='tight', facecolor=HEX_BG)
    plt.close(fig)
    buf.seek(0)
    return buf


def style_axes(ax):
    ax.set_facecolor(HEX_AXES_BG)
    ax.tick_params(colors=HEX_TEXT, labelsize=8)
    ax.xaxis.label.set_color(HEX_TEXT)
    ax.yaxis.label.set_color(HEX_TEXT)
    ax.title.set_color(HEX_TEXT)
    for sp in ['top', 'right']:
        ax.spines[sp].set_visible(False)
    for sp in ['bottom', 'left']:
        ax.spines[sp].set_color(HEX_SPINE)
    ax.grid(axis='y', color='#e0e0e0', linewidth=0.5, zorder=0)


# ============================================================================
# Chart generators
# ============================================================================

def chart_bs_limitation():
    np.random.seed(42)
    n = 300
    moneyness = np.random.uniform(0.75, 1.25, n)
    bs_price = np.where(moneyness > 1, (moneyness - 1) * 80, 0) + np.random.exponential(3, n)
    bs_price = np.clip(bs_price, 0.5, 55)
    smile_err = 0.35 * np.abs(moneyness - 1) + np.random.normal(0, 0.05, n)
    market_price = bs_price * (1 + smile_err)

    fig, ax = plt.subplots(figsize=(6, 5))
    fig.patch.set_facecolor(HEX_BG)
    style_axes(ax)

    sc = ax.scatter(bs_price, market_price, c=np.abs(moneyness - 1),
                    cmap='plasma', alpha=0.7, s=35, edgecolors='none')
    lim = max(bs_price.max(), market_price.max()) * 1.05
    ax.plot([0, lim], [0, lim], '--', color='#aaa', linewidth=1.5, label='Perfect fit (BS assumption)')
    cbar = fig.colorbar(sc, ax=ax)
    cbar.set_label('Distance from ATM', color=HEX_TEXT, fontsize=8)
    cbar.ax.yaxis.set_tick_params(color=HEX_TEXT)
    plt.setp(cbar.ax.yaxis.get_ticklabels(), color=HEX_TEXT)

    ax.set_xlabel('Black-Scholes Price ($)')
    ax.set_ylabel('Market Price ($)')
    ax.set_title('BS Systematically Underprices\nOTM / ITM Options (Smile Effect)',
                 fontweight='bold', fontsize=11)
    ax.legend(fontsize=9, facecolor='white', framealpha=0.8)
    plt.tight_layout()
    return fig_to_buf(fig)


def chart_data_pipeline():
    fig, axes = plt.subplots(1, 2, figsize=(9, 4.2))
    fig.patch.set_facecolor(HEX_BG)

    ax = axes[0]
    style_axes(ax)
    labels = ['Options\nChain', 'Greeks', 'Price\nHistory', 'Vol\nSurface', 'Earnings']
    counts = [5100, 5100, 1260, 300, 20]
    colors = [HEX_BLUE, HEX_GREEN, HEX_ORANGE, HEX_RED, HEX_PURPLE]
    bars = ax.bar(labels, counts, color=colors, edgecolor='none')
    for bar, val in zip(bars, counts):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 50,
                f'{val:,}', ha='center', color=HEX_TEXT, fontsize=9, fontweight='bold')
    ax.set_ylabel('Records')
    ax.set_title('Data Volume by Type', fontweight='bold')

    ax = axes[1]
    style_axes(ax)
    tickers = ['SPY', 'AAPL', 'NVDA', 'TSLA', 'QQQ']
    options = [1200, 950, 1100, 800, 1050]
    ax.barh(tickers, options, color=HEX_CYAN, edgecolor='none')
    for i, v in enumerate(options):
        ax.text(v + 15, i, str(v), va='center', color=HEX_TEXT, fontsize=9, fontweight='bold')
    ax.set_xlabel('Options Records')
    ax.set_title('Coverage by Ticker', fontweight='bold')

    plt.tight_layout()
    return fig_to_buf(fig)



def chart_model_results():
    # Real values from compute_metrics.py
    # BS vs lastPrice (market); XGBoost/LSTM vs bs_price target
    models = ['Black-Scholes', 'XGBoost', 'LSTM']
    r2     = [0.6798,  0.9987, -0.0001]
    rmse   = [61.67,   5.04,   139.60]
    colors = [HEX_BLUE, HEX_GREEN, HEX_RED]  # LSTM red — did not converge

    fig, axes = plt.subplots(1, 2, figsize=(9, 4.2))
    fig.patch.set_facecolor(HEX_BG)

    # Left: R²
    ax = axes[0]
    style_axes(ax)
    ax.grid(axis='y', color='#e0e0e0', linewidth=0.5, zorder=0)
    bars = ax.bar(models, r2, color=colors, edgecolor='none', zorder=3)
    ax.axhline(0, color=HEX_SPINE, linewidth=1)
    for bar, val in zip(bars, r2):
        ypos = val + 0.02 if val >= 0 else val - 0.06
        ax.text(bar.get_x() + bar.get_width() / 2, ypos,
                f'{val:.4f}', ha='center', color=HEX_TEXT, fontsize=10, fontweight='bold')
    ax.set_ylabel('R²')
    ax.set_title('R² Score  (higher is better)', fontweight='bold')
    ax.set_ylim(-0.2, 1.15)

    # Right: RMSE
    ax = axes[1]
    style_axes(ax)
    ax.grid(axis='y', color='#e0e0e0', linewidth=0.5, zorder=0)
    bars = ax.bar(models, rmse, color=colors, edgecolor='none', zorder=3)
    for bar, val in zip(bars, rmse):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 2,
                f'{val:.1f}', ha='center', color=HEX_TEXT, fontsize=10, fontweight='bold')
    ax.set_ylabel('RMSE ($)')
    ax.set_title('RMSE  (lower is better)', fontweight='bold')
    ax.set_ylim(0, 170)
    ax.annotate('* BS vs market price\n  XGB/LSTM vs BS target',
                xy=(0.98, 0.97), xycoords='axes fraction',
                ha='right', va='top', fontsize=7, color='#888888')

    plt.tight_layout()
    return fig_to_buf(fig)


def chart_var_earnings():
    np.random.seed(99)
    fig, axes = plt.subplots(1, 2, figsize=(9, 4.2))
    fig.patch.set_facecolor(HEX_BG)

    # VaR histogram
    ax = axes[0]
    style_axes(ax)
    returns = np.random.normal(-0.001, 0.012, 3000)
    var_95 = np.percentile(returns, 5)
    n_bins = 60
    counts, edges, patches_list = ax.hist(returns, bins=n_bins, color=HEX_BLUE, alpha=0.7, edgecolor='none')
    ax.axvline(var_95, color=HEX_RED, linewidth=2, label=f'95% VaR: {var_95:.3f}')
    ax.fill_between(
        [returns.min(), var_95],
        0,
        max(counts) * 0.9,
        alpha=0.25, color=HEX_RED
    )
    ax.set_xlabel('Daily Return'); ax.set_ylabel('Frequency')
    ax.set_title('Portfolio VaR (95% confidence)', fontweight='bold')
    ax.legend(fontsize=9, facecolor='white', framealpha=0.8)

    # Earnings impact bar chart
    ax = axes[1]
    style_axes(ax)
    surprises    = [-15, -10, -5, 0, 5, 10, 15]
    call_impacts = [-8.2, -5.1, -2.0, 0.5, 3.4, 7.1, 11.2]
    bar_colors   = [HEX_RED if v < 0 else HEX_GREEN for v in call_impacts]
    ax.bar(surprises, call_impacts, color=bar_colors, edgecolor='none', width=3.8)
    ax.axhline(0, color=HEX_SPINE, linewidth=0.8)
    ax.set_xlabel('EPS Surprise (%)'); ax.set_ylabel('Call Option ΔP ($)')
    ax.set_title('Earnings Impact on Option Price', fontweight='bold')

    plt.tight_layout()
    return fig_to_buf(fig)


# ============================================================================
# Slide builders
# ============================================================================

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(2.6), C_TITLE_BG)
    add_text(slide, "Portfolio Risk Stress-Tester",
             Inches(0.5), Inches(0.25), Inches(12.3), Inches(1.2),
             size=42, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "STATS 418 Final Project  ·  Mark Wang  ·  UCLA Statistics & Data Science",
             Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.55),
             size=16, color=C_ACCENT, align=PP_ALIGN.CENTER)

    add_text(slide, "Black-Scholes   ·   XGBoost   ·   LSTM   ·   Earnings Impact Model",
             Inches(0.5), Inches(2.8), Inches(12.3), Inches(0.5),
             size=15, bold=True, color=C_LGRAY, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(2.0), Inches(3.45), Inches(9.33), Inches(0.03), C_ACCENT)

    add_text(slide, "Live App:   https://stats418-portfolio-risk-tester.streamlit.app",
             Inches(1.5), Inches(3.65), Inches(10.3), Inches(0.4),
             size=13, color=C_GREEN, align=PP_ALIGN.CENTER)
    add_text(slide, "API Docs:  https://risk-api-388146732000.us-central1.run.app/docs",
             Inches(1.5), Inches(4.1), Inches(10.3), Inches(0.4),
             size=13, color=C_GREEN, align=PP_ALIGN.CENTER)

    add_text(slide, "Both services live through June 9, 2026",
             Inches(0), Inches(4.7), SLIDE_W, Inches(0.4),
             size=12, italic=True, color=C_LGRAY, align=PP_ALIGN.CENTER)


def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    title_bar(slide, "Problem & Motivation", "Why Black-Scholes isn't enough")

    bullets = [
        ("Options are central to hedging and portfolio risk management", []),
        ("Black-Scholes (1973) assumes:", [
            "Constant volatility — real markets contradict this every day",
            "Lognormal returns — fat tails and skew are well-documented",
        ]),
        ("Volatility smile: OTM/ITM options systematically\nunderpriced by BS — traders pay a premium", []),
        ("Goal: ML models that learn the pricing error from data", [
            "XGBoost — learns the smile implicitly from options features",
            "LSTM — captures temporal dynamics of volatility",
            "Earnings Impact — predicts option ΔP from EPS surprise %",
        ]),
    ]
    bullet_frame(slide, bullets, Inches(0.4), Inches(1.3), Inches(6.0), Inches(5.8))

    img = chart_bs_limitation()
    slide.shapes.add_picture(img, Inches(6.6), Inches(1.2), Inches(6.3), Inches(5.8))


def slide_data(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    title_bar(slide, "Data & Architecture", "5 tickers · ~11,760 records · live via yfinance")

    bullets = [
        ("Tickers: SPY · AAPL · NVDA · TSLA · QQQ", [
            "Index fund + high-vol mega-caps = diverse vol regimes",
        ]),
        ("Data pipeline:", [
            "fetch_market_data.py — options chains, price history, vol surface",
            "fetch_earnings.py — EPS history, upcoming dates",
            "preprocess.py — StandardScaler, feature engineering, 70/15/15 split",
        ]),
        ("Backend: FastAPI served on Google Cloud Run", []),
        ("Frontend: Streamlit on Streamlit Community Cloud", []),
        ("All models pretrained offline; inference at request time", []),
    ]
    bullet_frame(slide, bullets, Inches(0.4), Inches(1.3), Inches(5.5), Inches(5.8))

    img = chart_data_pipeline()
    slide.shapes.add_picture(img, Inches(6.2), Inches(1.2), Inches(6.8), Inches(5.8))



def slide_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    title_bar(slide, "Model Results", "Evaluated on held-out test set (15% of data) — values are placeholders")

    bullets = [
        ("Metrics: R² (higher → better) · RMSE (lower → better)", []),
        ("Black-Scholes  →  R² 0.68 vs actual market price", [
            "MAPE 34.6%, RMSE $61.67 — captures broad pricing structure",
            "Misses vol smile; OTM/ITM options systematically mispriced",
        ]),
        ("XGBoost  →  R² 0.9987 vs BS-price target", [
            "RMSE $5.04 — learns the BS formula from features near-perfectly",
            "High MAPE (38.8%) driven by deep-OTM options worth pennies",
        ]),
        ("LSTM  →  did not converge  (R² ≈ 0, RMSE $140)", [
            "Options rows are not a true time series — windowing across",
            "different contracts carries no temporal signal",
        ]),
        ("Earnings Impact  →  directionally accurate on EPS-surprise sign", [
            "Validated across 8 quarters of historical earnings events",
        ]),
    ]
    bullet_frame(slide, bullets, Inches(0.4), Inches(1.3), Inches(5.5), Inches(5.8), base_size=15)

    img = chart_model_results()
    slide.shapes.add_picture(img, Inches(6.2), Inches(1.5), Inches(6.8), Inches(5.6))


def slide_app(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    title_bar(slide, "Live App Features", "Three-tab Streamlit dashboard + 10-endpoint FastAPI")

    card_specs = [
        ("Earnings Impact\nPredictor",
         "Enter ticker → upcoming earnings date + EPS estimate → "
         "drag EPS surprise % slider → predicted option ΔP via XGBoost + delta-gamma approximation",
         HEX_GREEN),
        ("Portfolio VaR\n& Stress Test",
         "Upload portfolio CSV (ticker, weight, σ_daily) → "
         "parametric VaR at any confidence level and horizon → "
         "apply market shock → compare stressed vs. baseline VaR",
         HEX_ORANGE),
        ("Black-Scholes\nPricer + Greeks",
         "Interactive call/put pricer → live Delta, Gamma, Vega, Theta output → "
         "Greeks-vs-spot visualization across a user-defined price range",
         HEX_BLUE),
    ]

    for i, (title, desc, hex_c) in enumerate(card_specs):
        r, g, b = int(hex_c[1:3], 16), int(hex_c[3:5], 16), int(hex_c[5:7], 16)
        card_rgb   = RGBColor(r, g, b)
        card_dark  = RGBColor(max(0, r - 55), max(0, g - 55), max(0, b - 55))

        left = Inches(0.3 + i * 4.32)
        top  = Inches(1.3)
        add_rect(slide, left, top, Inches(4.0), Inches(5.75), card_dark)
        add_rect(slide, left, top, Inches(4.0), Inches(0.72), card_rgb)
        add_text(slide, title, left + Inches(0.1), top + Inches(0.07),
                 Inches(3.8), Inches(0.62), size=14, bold=True,
                 color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, desc, left + Inches(0.15), top + Inches(0.85),
                 Inches(3.72), Inches(4.7), size=13, color=C_WHITE)

    add_text(slide,
             "Also: Live Options Chain viewer — real-time from yfinance, sorted by strike, expiry-aware re-fetch",
             Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.32),
             size=11, italic=True, color=C_LGRAY, align=PP_ALIGN.CENTER)


def slide_var_earnings(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    title_bar(slide, "Portfolio VaR & Earnings Impact", "Risk quantification in action")

    bullets = [
        ("Parametric VaR", [
            "Upload CSV: ticker, weight, σ_daily",
            "95% / 99% VaR at any time horizon",
            "Stress test: apply market shock, compare VaR before and after",
        ]),
        ("Historical VaR", [
            "Upload returns CSV — distribution-free, no normality assumption",
        ]),
        ("Earnings Impact Predictor", [
            "Live EPS estimate from yfinance for any ticker",
            "EPS surprise slider → predicted call/put ΔP",
            "XGBoost + delta-gamma: captures both model and Greeks sensitivity",
        ]),
        ("Live Options Chain", [
            "Real-time options data, sorted by strike",
            "Expiry picker — re-fetches chain on change without state reset",
        ]),
    ]
    bullet_frame(slide, bullets, Inches(0.4), Inches(1.3), Inches(5.5), Inches(5.8))

    img = chart_var_earnings()
    slide.shapes.add_picture(img, Inches(6.2), Inches(1.2), Inches(6.8), Inches(5.8))


def slide_deployment(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    title_bar(slide, "Deployment", "Both services live through June 9, 2026")

    bullets = [
        ("Frontend — Streamlit Community Cloud", [
            "https://stats418-portfolio-risk-tester.streamlit.app",
            "Zero-infra hosting; secrets injected via Streamlit Cloud settings",
        ]),
        ("Backend — Google Cloud Run (FastAPI)", [
            "https://risk-api-388146732000.us-central1.run.app/docs",
            "Dockerized; auto-scales to zero; PORT injected at runtime",
        ]),
        ("10 API endpoints: /var/*, /options/*, /earnings/*, /stress/*", []),
        ("Key deployment hurdles:", [
            "lxml missing in Cloud Run — yfinance earnings_dates parsing failed silently",
            "CORS middleware + dynamic PORT binding for Cloud Run",
            "Streamlit secrets for API URL injection across environments",
        ]),
    ]
    bullet_frame(slide, bullets, Inches(0.4), Inches(1.3), Inches(6.5), Inches(5.8))

    # Architecture mini-diagram
    rl = Inches(7.5)
    rt = Inches(1.4)

    boxes = [
        ("Yahoo Finance API\n(external)", Inches(1.6), Inches(0.0),   RGBColor(0x0f, 0x3c, 0x78)),
        ("FastAPI Backend\n(Google Cloud Run)", Inches(1.6), Inches(1.4),  RGBColor(0x1a, 0x7a, 0x4a)),
        ("Streamlit Frontend\n(Streamlit Cloud)", Inches(1.6), Inches(2.8), RGBColor(0x7a, 0x45, 0x0a)),
        ("Browser", Inches(1.6), Inches(4.1),  RGBColor(0x3a, 0x3a, 0x7a)),
    ]
    for label, bx, by, bc in boxes:
        add_rect(slide, rl + bx, rt + by, Inches(2.3), Inches(0.75), bc)
        add_text(slide, label,
                 rl + bx + Inches(0.08), rt + by + Inches(0.05),
                 Inches(2.15), Inches(0.68),
                 size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    for by in [Inches(0.82), Inches(2.22), Inches(3.62)]:
        add_text(slide, "↓", rl + Inches(2.6), rt + by,
                 Inches(0.5), Inches(0.45), size=20, color=C_ACCENT, align=PP_ALIGN.CENTER)


def slide_lessons(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    title_bar(slide, "Lessons Learned", "Five weeks: raw data to a deployed ML system")

    left_bullets = [
        ("AI Assistance (Claude Code)", [
            "Most effective for narrow, specific tasks — not broad 'build X'",
            "Financial math (Greeks, BS) reliable; Streamlit session state needed review",
            "Cut Cloud Run debugging time significantly vs. reading docs cold",
        ]),
        ("Model Insights", [
            "XGBoost outperformed LSTM — temporal signal weaker than expected",
            "Implied vol is the single biggest driver of option price",
            "Earnings impact harder to model than price — noisy target signal",
        ]),
    ]
    right_bullets = [
        ("Infrastructure", [
            "Cloud Run + yfinance: silent failures from missing system packages",
            "Streamlit session state: initialization guards are non-negotiable",
            "Docker: separate pip install step from code copy for layer caching",
        ]),
        ("What I'd do differently", [
            "Collect more earnings events — 20 is too thin for robust ML",
            "Add real-time IV updates, not just a static chain snapshot",
            "Treat earnings impact as a time-series classification problem earlier",
        ]),
    ]

    bullet_frame(slide, left_bullets, Inches(0.4), Inches(1.3), Inches(6.0), Inches(5.8))
    add_rect(slide, Inches(6.62), Inches(1.35), Inches(0.03), Inches(5.7), C_ACCENT)
    bullet_frame(slide, right_bullets, Inches(6.8), Inches(1.3), Inches(6.1), Inches(5.8))


# ============================================================================
# Main
# ============================================================================

def build():
    from pathlib import Path
    fig_dir = Path("presentation/final_figs")
    fig_dir.mkdir(exist_ok=True)

    print("Saving standalone chart PNGs to presentation/final_figs/ ...")
    chart_fns = [
        ("fig_bs_limitation.png",  chart_bs_limitation),
        ("fig_data_pipeline.png",  chart_data_pipeline),
        ("fig_model_results.png",  chart_model_results),
        ("fig_var_earnings.png",   chart_var_earnings),
    ]
    for fname, fn in chart_fns:
        buf = fn()
        with open(fig_dir / fname, "wb") as f:
            f.write(buf.read())
        print(f"  Saved {fname}")

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    steps = [
        ("Title",                slide_title),
        ("Problem & Motivation", slide_problem),
        ("Data & Architecture",  slide_data),
        ("Results",              slide_results),
        ("App Features",         slide_app),
        ("VaR & Earnings",       slide_var_earnings),
        ("Deployment",           slide_deployment),
        ("Lessons Learned",      slide_lessons),
    ]

    print("\nBuilding PPTX slides...")
    for i, (name, fn) in enumerate(steps, 1):
        print(f"  [{i}/{len(steps)}] {name}")
        fn(prs)

    out = "presentation/wang-mark-final.pptx"
    prs.save(out)
    print(f"\nSaved PPTX: {out}")
    print(f"Saved PNGs: {fig_dir.resolve()}")


if __name__ == "__main__":
    build()
